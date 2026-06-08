#!/usr/bin/env python3
# 공식 PPTX 템플릿(슬라이드 2)에 연구 내용을 채워 완성본을 만든다.
# 폰트=맑은 고딕(템플릿 동일), 색=2C2C2C, 박스는 내용 영역에 맞게 리사이즈.
import copy
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

SRC = "assets/official_template.pptx"
OUT = "poster_kca_template.pptx"
FIG = "assets/fig"
FONT = "맑은 고딕"
TXT = RGBColor(0x2C, 0x2C, 0x2C)

prs = Presentation(SRC)
slide = prs.slides[1]  # slide 2 = template

# 섹션 바·라벨 재배치 (배경 LaTeX판과 동일): Y 이동 + 가운데로 모으기 + 라벨 수직 중앙
def _find(name):
    for sp in slide.shapes:
        if sp.name == name: return sp
    return None
def shift_x(name, dx_mm):
    sp = _find(name);
    if sp is not None: sp.left = sp.left + Mm(dx_mm)
def set_y(name, y_mm):
    sp = _find(name);
    if sp is not None: sp.top = Mm(y_mm)
def center_label(lbl, bar_top_mm, dx_mm, bar_h_mm=34.3):
    sp = _find(lbl)
    if sp is None: return
    sp.left = sp.left + Mm(dx_mm); sp.top = Mm(bar_top_mm); sp.height = Mm(bar_h_mm)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0; tf.margin_bottom = 0
L, R = 6, -10  # 좌측 +6, 우측 -10 → 가운데로 (gutter 축소)
for nm in ["Picture 5", "그룹 35", "그룹 39", "Picture 11", "Picture 15"]: shift_x(nm, L)
for nm in ["Picture 13", "Picture 17"]: shift_x(nm, R)
set_y("그룹 39", 445); set_y("Picture 11", 565); set_y("Picture 15", 760)
center_label("TextBox 8", 179.8, L); center_label("TextBox 12", 565, L); center_label("TextBox 16", 760, L)
center_label("TextBox 14", 179.8, R); center_label("TextBox 18", 888.7, R)

def find(sub):
    for sp in slide.shapes:
        if sp.has_text_frame and sub in sp.text_frame.text:
            return sp
    raise KeyError(sub)

def style_run(r, size, bold=False, color=TXT):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    # set East-Asian typeface so Korean uses Malgun Gothic too
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)

def setbox(sp, x, y, w, h):
    sp.left, sp.top, sp.width, sp.height = Mm(x), Mm(y), Mm(w), Mm(h)
    tf = sp.text_frame
    tf.word_wrap = True
    try: tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception: pass
    tf.vertical_anchor = MSO_ANCHOR.TOP

def fill_para(sp, paras, size, bold_lead=False, bullet=False):
    """paras: list of (text) or (lead, rest) tuples."""
    tf = sp.text_frame
    tf.clear()
    for i, item in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = lead; style_run(r1, size, bold=True)
            r2 = p.add_run(); r2.text = rest; style_run(r2, size)
        else:
            r = p.add_run(); r.text = item; style_run(r, size)

# ---------------- TITLE / AUTHOR ----------------
t = find("논문명 기재")
setbox(t, 70, 60, 701, 70)
tf = t.text_frame; tf.clear()
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "한국어 공개 SNS에 기록된 상담 경험의 기술적 특징"; style_run(r, 60, bold=True)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "LLM 보조 정제 절차로 분석한 사례 보고"; style_run(r2, 30)

a = find("이름(소속)")
setbox(a, 70, 126, 701, 24)
tf = a.text_frame; tf.clear()
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "이윤정 (서울대학교)"; style_run(r, 34, bold=True)

# ---------------- LEFT COLUMN ----------------
seoron = ("상담 경험 연구는 부정 경험을 충분히 담아내지 못해 왔다. 그 누락은 부정 경험을 가진 사람이 "
"면접 표본에 들어오지 않는 모집단 차원보다, 같은 사람이라도 면접 환경에서는 그 경험을 보고하지 않는 "
"보고 환경 차원에서 강하게 일어난다(Mohr, 1995; Krumpal, 2013; 문보경·장성숙, 2001). 부정 경험을 가진 "
"내담자일수록 추적 면접에 잘 응하지 않고, 대면 맥락에서는 사회적 바람직성 편향이 작용하기 때문이다. "
"면접의 한계가 보고 환경에서 비롯된 것이라면, 환경을 바꿀 때 같은 경험이 다른 형태로 기록될 수 있다. "
"익명을 전제로 자기 경험을 적는 SNS 글이 그 사례다. Steinbrenner 외(2025)는 독일어권 온라인 포럼에서 "
"부정적 심리치료 경험의 주제 구조를 추출해, 공개 텍스트가 분포 수준의 자료원이 됨을 보였다. 그러나 "
"한국어 회고 표현과 상담 문화의 특수성 때문에 영어·독일어권 결과를 그대로 옮기기 어렵고, 한국어 자료를 "
"직접 분석한 연구는 아직 드물다.")
sp = find("연구의 필요성을 작성"); setbox(sp, 58, 250, 345, 190); fill_para(sp, [seoron], 27)

sp = find("연구목적과 연구문제를 작성"); setbox(sp, 61, 470, 345, 90)
fill_para(sp, [("목적 1.  ", "한국어 공개 SNS 상담 경험 자료에 적용할 LLM 보조 정제 절차를 정리하고, 본문 코퍼스를 구성한다."),
               ("목적 2.  ", "본문 코퍼스의 경험 극성과 사건 언어 양상을 분포 수준에서 기술하고, 면접 기반 선행 연구와 비교한다.")], 22)

iron1 = ("본 연구의 분석 틀은 두 층으로 이루어진다. 이론 고정층은 작업동맹이다. Bordin(1979)은 작업동맹을 "
"유대·과제·목표 세 축의 협력으로 보았다. 유대는 상담자와 내담자의 정서적 결속, 과제는 회기에서 합의해 "
"수행하는 활동, 목표는 함께 이루려는 변화를 가리킨다. 이 세 축은 짧은 회고문에서도 부정 경험과 긍정 경험을 "
"다른 종류의 사건으로 구분하게 해 준다.")
iron2 = ("탐색 고정층은 부정 경험 사건 언어에서 끌어왔다. Mohr(1995)는 부정적 성과를 식별 가능한 사건 단위로 "
"분리해 보고할 것을 권고했고, 문보경·장성숙(2001)은 내담자 불만을 '원치 않는 반응', '상담자의 요구', "
"'상처 경험' 등으로 범주화했다. 본 연구는 이 사건 언어를 살려 여섯 과정 코드를 두었다. 부정 쪽은 "
"강제성·무시당함·상처, 긍정 쪽은 이해받음·안도·성장이다. 각 코드는 글마다 등장 여부를 이진값으로 본다.")
sp = find("이론적 배경을 작성"); setbox(sp, 55, 610, 345, 145); fill_para(sp, [iron1, iron2], 23)

def caption(x, y, w, text, size=14):
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(14)); tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.italic = True; r.font.color.rgb = TXT
    return tb
bangbeop = ("자료는 X 공개 한국어 글에서 모았다(사용자명·외부 식별자 자동 비식별). 질의는 정밀·확장·맥락 점검·"
"누적 확장의 네 층으로 설계했고, 회수 게시물은 2010년 9월~2026년 4월(약 15년 7개월)에 걸쳐 있다.")
sp = find("연구방법에 대해 작성"); setbox(sp, 55, 805, 345, 50); fill_para(sp, [bangbeop], 22)
# 전체 폭 4층 파이프라인 + 그 아래 코퍼스 퍼널 (원본 포스터와 동일 구성)
slide.shapes.add_picture(f"{FIG}/pipeline.png", Mm(56), Mm(856), width=Mm(338))
caption(56, 978, 338, "그림 1. LLM 보조 4층 정제 절차. 규칙 분류로 분량을 줄이고 경계 사례에만 LLM 판정(③)을 적용한다.")
slide.shapes.add_picture(f"{FIG}/corpus_funnel.png", Mm(118), Mm(1002), width=Mm(205))
caption(118, 1118, 205, "그림 2. 원자료 3,652건 → 본문 코퍼스 361건(9.9%). 관련성 유지본 479·엄격 확정본 451건 경유.")

# ---------------- RIGHT COLUMN ----------------
gi = ("본문 코퍼스 361건의 경험 극성은 부정 117·혼합 91·긍정 126·중립 27건이다. 정식 표본에서 과소표집되던 "
"부정 경험이 긍정과 엇비슷한 규모로 함께 들어왔다. 장면은 일반 정신건강 상담 268건과 청소년 학교상담 "
"93건(약 26%)으로 갈렸다.")
sp = find("연구결과를 작성"); setbox(sp, 449, 224, 338, 66)
fill_para(sp, [("기초 분포 — ", gi)], 24)
# overview figure
slide.shapes.add_picture(f"{FIG}/corpus_overview.png", Mm(453), Mm(290), width=Mm(330))

# 비교 textbox (new)
tb = slide.shapes.add_textbox(Mm(449), Mm(420), Mm(338), Mm(52))
fill_para(tb, [("면접 기반 선행 연구와의 비교 — ",
  "부정 117건에서 관찰된 사건 표지를 문보경·장성숙(2001)의 면접 범주와 나란히 놓았다. 면접 자료에서 "
  "사건 단위로 정리된 부정 경험의 표지가 익명 환경의 글에서도 적용 가능한 형태로 드러난다.")], 24)
tb.text_frame.word_wrap = True

# comparison TABLE (native)
rows = [
  ("측면", "면접 자료 (문보경·장성숙, 2001)", "본 SNS 자료"),
  ("자료원", "내담자 면접·개방형 질문지", "한국어 공개 X 경험 글"),
  ("분석 단위", "사례 단위 사건 회고", "게시물 단위 경험 글"),
  ("부정 표지", "원치 않는 반응, 해결책 부재, 상담자의 요구, 상처 경험", "강제성 0.350, 무시당함 0.470, 상처 0.231 (부정 117건 기준)"),
  ("분석 깊이", "사례별 깊이 있는 합의 부호화", "분포 수준의 탐색적 부호화"),
]
gf = slide.shapes.add_table(len(rows), 3, Mm(449), Mm(478), Mm(338), Mm(96))
tbl = gf.table
tbl.first_row = False; tbl.horz_banding = False  # 직접 색칠 (기본 스타일 줄무늬 끔)
# 기본 표 스타일 제거(No Style, No Grid) → 명시적 셀 색이 그대로 보이도록
_tblPr = tbl._tbl.tblPr
for _el in _tblPr.findall(qn('a:tableStyleId')): _tblPr.remove(_el)
_sid = _tblPr.makeelement(qn('a:tableStyleId'), {}); _sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
_tblPr.append(_sid)
tbl.columns[0].width = Mm(58); tbl.columns[1].width = Mm(142); tbl.columns[2].width = Mm(138)
NAVY = RGBColor(0x14, 0x28, 0x5A); LIGHT = RGBColor(0xEC, 0xEF, 0xF6); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Mm(2.5); cell.margin_right = Mm(2.5)
        cell.margin_top = Mm(1.5); cell.margin_bottom = Mm(1.5)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (LIGHT if ri % 2 == 1 else WHITE)
        tf = cell.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run(); r.text = val
        style_run(r, 17, bold=(ri == 0 or ci == 0), color=(WHITE if ri == 0 else TXT))

# 과정 코드 textbox (new)
gc = ("작업동맹 세 축은 극성에 따라 갈렸다. 부정 117건에서는 유대 63.2%·과제 69.2%·목표 59.0%가 부정 정렬, "
"긍정 126건에서는 유대 80.2%·과제 81.7%·목표 84.9%가 긍정 정렬이었다. 여섯 과정 코드의 부트스트랩 대비는 "
"신뢰구간이 모두 0을 가로지르지 않았다(예: 무시당함 -0.431 [-0.506, -0.350], 안도 +0.434 [0.344, 0.517]).")
tb2 = slide.shapes.add_textbox(Mm(449), Mm(582), Mm(338), Mm(80))
fill_para(tb2, [("과정 코드의 극성별 대비 — ", gc)], 24)
tb2.text_frame.word_wrap = True
# process figure
slide.shapes.add_picture(f"{FIG}/process_contrast.png", Mm(460), Mm(672), width=Mm(305))

# ---------------- V. 논의 ----------------
sp = find("결론 및 논의를 작성"); setbox(sp, 449, 936, 338, 200)
fill_para(sp, [
  ("함의:  ", "정식 표본에서 좀처럼 함께 모이지 않던 부정과 긍정 경험이 익명 환경에서는 거의 같은 비중으로 함께 "
   "기록되었다(부정 117·긍정 126건). 사건 단위의 경험 언어도 분포 수준에서 갈렸으며, 그 표지는 면접 자료를 "
   "분석한 문보경·장성숙(2001)과 닮았다. 다만 이 일치는 같은 모집단의 일치가 아니라, 본 연구의 코드가 그 "
   "표지를 출발점으로 삼은 데서 비롯한다."),
  ("한계:  ", "본 결과는 사람 합의 코딩층 없이 LLM 보조 부호화로만 얻었고, 익명 환경에서 자기 경험을 적는 "
   "사람이라는 표집 편향 위에서 관찰되었다."),
  ("전망:  ", "따라서 본 연구는 한국 상담 경험 일반에 대한 결론이 아닌 사례 보고이며, 후속 합의적 질적 연구와 "
   "면접·SNS 직접 비교의 출발 자료로 쓸 수 있다."),
], 24)

prs.save(OUT)
print("saved", OUT)
